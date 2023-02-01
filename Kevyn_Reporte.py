from io import BytesIO
import io
from math import isnan
import os
import numpy as np
from pptx.util import Inches
import pandas as pd
import streamlit as st
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
import requests
from pptx.chart.data import CategoryChartData
import plotly.graph_objects as go

def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell


def format_cpc_cells(n):
    return f"USD {'{0:.2f}'.format(n)}"


def format_cells(n):
    return "{0:.2f}".format(n)


def color_negative_red(value):
    color = "green"
    return "color: %s" % color


def highlight(s):
    if s.duration < 3:
        return ["background-color: yellow"] * len(s)
    else:
        return ["background-color: red"] * len(s)


# Funcion para generar archivo pptx
def generate_pptx(prs):
    print("-------")
    if facebook_csv is not None:
        fb_df = pd.read_csv(facebook_csv, header=None, sep=None, encoding = "ISO-8859-1")
        edad_sexo_start_row, edad_sexo_start_col = np.where(fb_df == "Edad y sexo")    
        edad_sexo_start_col = edad_sexo_start_col[0]
        edad_sexo_start_row = edad_sexo_start_row[0] + 1
        edad_sexo_df = fb_df.iloc[edad_sexo_start_row:edad_sexo_start_row+7, edad_sexo_start_col:edad_sexo_start_col+3]
        edad_sexo_row, edad_sexo_col = edad_sexo_df.shape
        for i in range(edad_sexo_row):
            for j in range(edad_sexo_col):
                if '%' in edad_sexo_df.iloc[i,j]:
                    edad_sexo_df.iloc[i,j] = float(edad_sexo_df.iloc[i,j].rstrip('%').replace(',', '.')) / 100
        hombre_count = sum(edad_sexo_df.iloc[1:7,2])
        mujer_count = sum(edad_sexo_df.iloc[1:7,1])
        hombre_series = edad_sexo_df.iloc[1:7, 2].values
        mujer_series =  edad_sexo_df.iloc[1:7, 1].values
        edad_sexo_categories = edad_sexo_df.iloc[1:7, 0].values
        pagina_start_row, pagina_start_col = np.where(fb_df == 'Principales páginas')
        pagina_start_col = pagina_start_col[0]
        pagina_start_row = pagina_start_row[0] + 1
        pagina_df = fb_df.iloc[pagina_start_row:pagina_start_row+11, pagina_start_col:pagina_start_col+2]
        pagina_row, pagina_col = pagina_df.shape
        print(pagina_df)
        for i in range(pagina_row):
            for j in range(pagina_col):
                if '%' in pagina_df.iloc[i,j]:
                    pagina_df.iloc[i,j] = float(pagina_df.iloc[i,j].rstrip('%').replace(',', '.')) / 100
        pagina_categories = pagina_df.iloc[1:11, 0].values
        pagina_series = pagina_df.iloc[1:11, 1].values
        for shape in prs.slides[9].shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                chart_title = shape.chart.chart_title.text_frame.text                        
                if 'Edad' in chart_title:
                        chart_data = CategoryChartData()
                        chart_data.categories = edad_sexo_categories
                        chart_data.add_series("Hombre", hombre_series)
                        chart_data.add_series("Mujer", mujer_series)
                        shape.chart.replace_data(chart_data)
                if 'Principales' in chart_title:
                        chart_data = CategoryChartData()
                        chart_data.categories = pagina_categories
                        chart_data.add_series("Paginas", pagina_series)
                        shape.chart.replace_data(chart_data)
                if chart_title == "Género":
                    if generos is not None:
                        chart_data = CategoryChartData()
                        chart_data.categories = ['Hombre', 'Mujer']
                        chart_data.add_series("Género", [hombre_count, mujer_count])
                        shape.chart.replace_data(chart_data)


    for shape in prs.slides[6].shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and shape.text == "title":
            shape.text = ""
            frame2 = shape.text_frame.paragraphs[0]
            frame2.alignment = PP_ALIGN.CENTER
            run2 = frame2.add_run()
            run2.text = name
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            chart_title = shape.chart.chart_title.text_frame.text
            if chart_title == "Género":
                genero_categories = []
                genero_series = ()
                if generos is not None:
                    generos_df = pd.read_csv(generos, header=None)
                    genero_categories = generos_df.iloc[0].values[1:]
                    genero_series = generos_df.iloc[1].values[1:]
                    genero_sum = sum(map(int, genero_series))
                    genero_series = [int(x) / genero_sum for x in genero_series]
                    chart_data = CategoryChartData()
                    chart_data.categories = genero_categories
                    chart_data.add_series("Género", genero_series)
                    shape.chart.replace_data(chart_data)
            if chart_title == "Etapa familiar":
                etapa_categories = []
                etapa_series = ()
                if etapas is not None:
                    etapas_df = pd.read_csv(etapas, header=None)
                    etapa_categories = etapas_df.iloc[0].values[1:]
                    etapa_series = etapas_df.iloc[1].values[1:]
                    etapa_sum = sum(map(int, etapa_series))
                    etapa_series = [int(x) / etapa_sum for x in etapa_series]
                    chart_data = CategoryChartData()
                    chart_data.categories = etapa_categories
                    chart_data.add_series("Etapa Familiar", etapa_series)
                    shape.chart.replace_data(chart_data)
            if chart_title == "Edad":
                edad_categories = []
                edad_series = ()
                if edades is not None:
                    df = pd.read_csv(edades, header=None)
                    edad_categories = df.iloc[:, 0].values[1:]
                    edad_series = df.iloc[:, 1].values[1:]
                    edad_sum = sum(map(int, edad_series))
                    edad_series = [int(x) / edad_sum for x in edad_series]
                    chart_data = CategoryChartData()
                    chart_data.categories = edad_categories
                    chart_data.add_series("Edad", edad_series)
                    shape.chart.replace_data(chart_data)
            if chart_title == "Ocupaciones":
                ocupacion_categories = []
                ocupacion_series = ()
                if ocupaciones is not None:
                    ocupaciones_df = pd.read_csv(ocupaciones, header=None)
                    ocupacion_categories = ocupaciones_df.iloc[0].values[1:]
                    ocupacion_series = ocupaciones_df.iloc[1].values[1:]
                    ocupacion_sum = sum(map(int, ocupacion_series))
                    ocupacion_series = [
                        int(x) / ocupacion_sum for x in ocupacion_series
                    ]
                    chart_data = CategoryChartData()
                    chart_data.categories = ocupacion_categories
                    chart_data.add_series("Ocupaciones", ocupacion_series)
                    shape.chart.replace_data(chart_data)
            if chart_title == "Intereses":
                interes_categories = []
                interes_series = ()
                if intereses is not None:
                    intereses_df = pd.read_csv(intereses, header=None)
                    interes_categories = intereses_df.iloc[0].values[1:]
                    interes_series = intereses_df.iloc[1].values[1:]
                    interes_sum = sum(map(int, interes_series))
                    interes_series = [int(x) / interes_sum for x in interes_series]
                    chart_data = CategoryChartData()
                    chart_data.categories = interes_categories
                    chart_data.add_series("Intereses", interes_series)
                    shape.chart.replace_data(chart_data)
    
    # for shape in prs.slides[8].shapes:
    #     print(shape.shape_type)
    k = 0.1
    indexpd = []
    indexC = -1
    dataTableName = []
    dataTableAverage = []
    for file in keyword_files:
        indexC += 1
        if keyword_files is not None:
            keyword_df = pd.read_csv(file, sep=None, encoding = "ISO-8859-1")           
            row, col = keyword_df.shape
            promedio = 0.0

            col_formats = {"Competition": ".2%"}
            plot_df = keyword_df.reset_index()
            font_colours_df = pd.DataFrame(
                "black",  # Set default font colour
                index=plot_df.index,
                columns=plot_df.columns,
            )
            colors = []
            colors = ["rgb(255, 255, 255)" for i in range(row)]

            for i in range(col):
                for j in range(row):
                    if j == 0 or j == 1 or j == 2:
                        colors[j] = "rgb(252, 248, 3)"
                    # if j + 1 < row:
                        # table.cell(j + 1, i).text = str(keyword_df.iloc[j + 1, i])
                    if i == 1:
                        promedio += float(keyword_df.iloc[j, i])

            dataTableName.append(os.path.splitext(file.name)[0])
            dataTableAverage.append(str(round((promedio/row), 2)))
            indexpd.append(indexC)
            data = {}
            for col_name in keyword_df.columns:
                data[col_name] = keyword_df[col_name].to_numpy()
            data["Color"] = colors
            df = pd.DataFrame(data)
            df.fillna('', inplace=True)

            fig = go.Figure(
                data=[
                    go.Table(
                        columnwidth=[0.3, 0.1, 0.1, 0.1, 0.1],
                        header=dict(
                            values=[col for col in df if col != "Color"],
                            line_color="black",
                            fill_color="white",
                            align="center",
                            font=dict(color="black", size=12),
                        ),
                        cells=dict(
                            values=[df[col].values for col in df if col != "Color"],
                            line_color=["rgb(0, 0, 0)"] * len(colors),
                            fill_color=[df.Color],
                            align=["center"],
                            font=dict(color="black", size=11),
                        ),
                    )
                ],
            )
            fig.update_layout(
                height=800, margin={"t": 1, "b": 1, "r": 1, "l": 1}, width=600
            )
            # fig.show()
            fig.write_image("plot1.png")

            shapes = prs.slides[8].shapes
            


            # table = shapes.add_table(row, col, left, top, width, height).table

            # for i in range(col):
            #     table.columns[i].width = Inches(1.0)
            #     table.cell(0, i).text = encabezados[i]
            # print(col)
            # for i in range(col):
            #     for j in range(row):
            #         if j + 1 < row:
            #             # table.cell(j + 1, i).text = str(keyword_df.iloc[j + 1, i])
            #             if i == 1:
            #                 promedio += float(keyword_df.iloc[j + 1, i])

            # print(promedio)

            shapes.add_picture("plot1.png", Inches(k), Inches(1))
            k += 6  

        
        # data2 = {}
    colors2 = []
    colors2 = ["rgb(255, 255, 255)" for i in range(indexC+1)]

    ex_dict = {
        'Categoria': dataTableName,
        'Búsqueda mensual promedio': dataTableAverage,
        'Color': colors2 
    }
    df2 = pd.DataFrame(ex_dict, columns=['Categoria','Búsqueda mensual promedio', 'Color'], index=indexpd)
    
    
    if len(keyword_files) > 0 :
        fig2 = go.Figure(
                data=[
                    go.Table(
                        columnwidth=[0.3, 0.1, 0.1, 0.1, 0.1],
                        header=dict(
                            values=['Categoria','Búsqueda mensual promedio'],
                            line_color="black",
                            fill_color="white",
                            align="center",
                            font=dict(color="black", size=12),
                        ),
                        cells=dict(
                            values=[df2[col].values for col in df2 if col != "Color"],
                            line_color=["rgb(0, 0, 0)"] * len(colors2),
                            fill_color=[df2.Color],
                            align=["center"],
                            font=dict(color="black", size=11),
                        ),
                    )
                ],
            )
        fig2.update_layout(
            height=300, margin={"t": 1, "b": 1, "r": 1, "l": 1}, width=400
        )
            # fig.show()
        fig2.write_image("plot2.png")

        shapes = prs.slides[7].shapes
        shapes.add_picture("plot2.png", Inches(1), Inches(3))

        # posInicio = 1
        # shapes2 = prs.slides[7].shapes
        # left = Inches(1.0)
        # top = Inches(4.0)
        # width = Inches(5.0)
        # height = Inches(2.0)
        # table2 = shapes2.add_table(len(keyword_files)+1, 2, left, top, width, height).table

        # # cell is a table cell
        # # set fill type to solid color first
        # table2.cell(0, 0).fill.solid()

        # # set foreground (fill) color to a specific RGB color
        # table2.cell(0, 0).fill.fore_color.rgb = RGBColor(0xFB, 0x8F, 0x00)
        # table2.columns[0].width = Inches(3.0)
        # table2.columns[1].width = Inches(3.0)
        # table2.cell(0, 0).text = 'Categoria'
        # table2.cell(0, 1).text = 'Busqueda mensual promedio'
        # for i in range(len(keyword_files)):
        #     table2.cell(posInicio,0).text = dataTableName[i]
        #     table2.cell(posInicio,1).text = dataTableAverage[i]
        #     posInicio += 1

    # row col
    # set column widths
    # table.columns[0].width = Inches(1.0)
    # table.columns[1].width = Inches(1.0)
    # table.columns[2].width = Inches(1.0)
    # table.columns[3].width = Inches(1.0)
    # table.columns[4].width = Inches(1.0)

    # write column headings
    # table.cell(0, 0).text = 'Keyword'
    # table.cell(0, 1).text = 'Búsqueda Promedio'
    # table.cell(0, 2).text = 'Impresiones'
    # table.cell(0, 3).text = 'Competencia'
    # table.cell(0, 4).text = 'CPC Promedio'

    # write body cells
    # table.cell(1, 0).text = 'Baz'
    # table.cell(1, 1).text = 'Qux'
    # for cell in iter_cells(table):
    #     for paragraph in cell.text_frame.paragraphs:
    #         for run in paragraph.runs:
    #             run.font.size = Pt(1)
    binary_output = BytesIO()
    prs.save(binary_output)
    return binary_output.getvalue()


st.set_page_config(layout="wide")
st.title("Reporte")

name = st.text_input(label="Nombre del proyecto", key="name")

edades = st.file_uploader("CSV Edades")
ocupaciones = st.file_uploader("CSV Ocupaciones")
etapas = st.file_uploader("CSV Etapas")
intereses = st.file_uploader("CSV Intereses")
generos = st.file_uploader("CSV Generos")

keyword_files = st.file_uploader("Google Ads", accept_multiple_files=True)

facebook_csv = st.file_uploader("Facebook CSV")


r = requests.get(
    "https://github.com/Neuronsinc/DataMiningReport/blob/main/template.pptx?raw=true"
)
prs = Presentation(io.BytesIO(r.content))


st.download_button(
    label="Descargar pptx",
    data=generate_pptx(prs),
    file_name="report.pptx",
    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
)
